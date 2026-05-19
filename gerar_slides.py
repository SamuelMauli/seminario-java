"""
Gera apresentacao do seminario de Concorrencia em Java.
Saida: slides/Seminario_Concorrencia_Java.pptx
"""
import os
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(OUT_DIR, "slides")
IMG_DIR = os.path.join(SLIDES_DIR, "img")
os.makedirs(IMG_DIR, exist_ok=True)

AZUL   = RGBColor(0x0B, 0x3D, 0x91)
LARANJA= RGBColor(0xE8, 0x7A, 0x1E)
CINZA  = RGBColor(0x44, 0x44, 0x44)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
FUNDO  = RGBColor(0xF6, 0xF8, 0xFC)

# ------------------- GRAFICOS -------------------
def grafico_race_condition():
    fig, ax = plt.subplots(figsize=(7, 4))
    cenarios = ["Sem sync\n(race)", "Com synchronized"]
    valores = [134289, 200000]
    cores = ["#E8531E", "#0B8A3E"]
    bars = ax.bar(cenarios, valores, color=cores)
    ax.axhline(200000, ls="--", color="gray", label="Esperado = 200.000")
    for b, v in zip(bars, valores):
        ax.text(b.get_x()+b.get_width()/2, v+3000, f"{v:,}", ha="center", fontsize=12, fontweight="bold")
    ax.set_ylabel("Valor final do contador")
    ax.set_title("Race Condition: 2 threads incrementando 100.000x")
    ax.legend()
    ax.set_ylim(0, 230000)
    plt.tight_layout()
    path = os.path.join(IMG_DIR, "race.png")
    plt.savefig(path, dpi=140); plt.close()
    return path

def grafico_ciclo_vida():
    fig, ax = plt.subplots(figsize=(8, 4.2))
    estados = ["NEW", "RUNNABLE", "BLOCKED", "WAITING", "TIMED_WAITING", "TERMINATED"]
    descricoes = ["criada", "executando", "esperando lock", "wait/join", "sleep/wait(ms)", "fim"]
    y = list(range(len(estados)))[::-1]
    ax.barh(y, [1]*len(estados), color=["#9aa", "#0B8A3E", "#E8531E", "#1E6FE8", "#E8C01E", "#666"])
    for i, (e, d) in enumerate(zip(estados, descricoes)):
        ax.text(0.5, y[i], f"{e} — {d}", ha="center", va="center", color="white", fontsize=12, fontweight="bold")
    ax.set_xlim(0, 1); ax.set_ylim(-0.5, len(estados)-0.5)
    ax.axis("off")
    ax.set_title("Ciclo de Vida de uma Thread (Thread.State)")
    plt.tight_layout()
    path = os.path.join(IMG_DIR, "ciclo.png")
    plt.savefig(path, dpi=140); plt.close()
    return path

def grafico_economia():
    fig, ax = plt.subplots(figsize=(7, 4))
    tecnicas = ["Imutável\n(final)", "volatile", "synchronized", "wait/notify", "ThreadLocal"]
    seguranca = [10, 4, 9, 9, 8]
    custo     = [1, 2, 6, 7, 3]
    x = range(len(tecnicas))
    ax.bar([i-0.18 for i in x], seguranca, width=0.36, color="#0B3D91", label="Segurança")
    ax.bar([i+0.18 for i in x], custo,     width=0.36, color="#E8531E", label="Custo / complexidade")
    ax.set_xticks(list(x)); ax.set_xticklabels(tecnicas)
    ax.set_ylim(0, 11)
    ax.set_title("Trade-off: Segurança vs Custo (0–10)")
    ax.legend()
    plt.tight_layout()
    path = os.path.join(IMG_DIR, "tradeoff.png")
    plt.savefig(path, dpi=140); plt.close()
    return path

def grafico_problemas():
    fig, ax = plt.subplots(figsize=(7, 4))
    problemas = ["Race\nCondition", "Deadlock", "Starvation", "Livelock"]
    impacto   = [9, 10, 6, 7]
    cores = ["#E8531E", "#B11E1E", "#E8C01E", "#1E6FE8"]
    bars = ax.bar(problemas, impacto, color=cores)
    for b, v in zip(bars, impacto):
        ax.text(b.get_x()+b.get_width()/2, v+0.2, str(v), ha="center", fontsize=12, fontweight="bold")
    ax.set_ylim(0, 12)
    ax.set_ylabel("Severidade típica (0–10)")
    ax.set_title("Problemas Clássicos de Concorrência")
    plt.tight_layout()
    path = os.path.join(IMG_DIR, "problemas.png")
    plt.savefig(path, dpi=140); plt.close()
    return path

# ------------------- HELPERS PPTX -------------------
def add_slide(prs, layout=6):
    return prs.slides.add_slide(prs.slide_layouts[layout])

def set_background(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)

def add_text(slide, text, left, top, width, height, size=18, bold=False, color=CINZA, align=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None: p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return tb

def add_title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.2), prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.5)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BRANCO

def add_bullets(slide, items, left, top, width, height, size=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = "• " + it
        r.font.size = Pt(size); r.font.color.rgb = CINZA
        p.space_after = Pt(6)

def add_code(slide, code, left, top, width, height, size=12):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.fill.solid();
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line if line else " "
        r.font.name = "Menlo"; r.font.size = Pt(size); r.font.color.rgb = RGBColor(0x1A,0x1A,0x2E)

# ------------------- MONTAGEM -------------------
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

# 1 — Capa
s = add_slide(prs); set_background(s, FUNDO)
add_text(s, "Concorrência em Java", Inches(0.6), Inches(1.6), Inches(12), Inches(1.2),
         size=54, bold=True, color=AZUL)
add_text(s, "Threads, Memória e Coordenação Básica", Inches(0.6), Inches(2.7), Inches(12), Inches(0.8),
         size=28, color=LARANJA)
add_text(s, "Equipe (ordem alfabética):", Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
         size=20, bold=True, color=CINZA)
equipe = "Lucas Gomes  •  Luiz Eduardo  •  Luiz Gabriel  •  Pedro Rossi  •  Renan Souza  •  Samuel Mauli  •  William Ferraz"
add_text(s, equipe, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6), size=18, color=CINZA)
add_text(s, "Seminário • 2026", Inches(0.6), Inches(6.6), Inches(12), Inches(0.5), size=16, color=AZUL)

# 2 — Agenda
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "Agenda")
add_bullets(s, [
    "1. Criação e controle de threads (Thread, Runnable)",
    "2. Métodos: start, sleep, join, yield, interrupt",
    "3. Daemon vs User threads + Ciclo de vida",
    "4. Sincronização: synchronized, volatile, final",
    "5. Monitores e wait/notify/notifyAll",
    "6. Java Memory Model (happens-before)",
    "7. Problemas: race, deadlock, starvation, livelock",
    "8. Contexto por thread: ThreadLocal e InheritableThreadLocal",
    "9. Demo do código + Conclusão",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5.5), size=22)

# 3 — Thread vs Runnable
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "1) Thread vs Runnable")
add_bullets(s, [
    "Thread = a unidade de execução da JVM",
    "Runnable = a TAREFA que a thread roda (preferido)",
    "start() inicia em paralelo; run() chamado direto roda na main",
], Inches(0.6), Inches(1.4), Inches(6.2), Inches(4))
add_code(s, """new Thread(() -> {
    System.out.println("rodando em "
        + Thread.currentThread().getName());
}, "worker").start();""",
    Inches(6.9), Inches(1.5), Inches(6), Inches(2.5), size=14)

# 4 — Métodos importantes
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "2) Métodos importantes")
add_bullets(s, [
    "start()        — inicia a thread",
    "sleep(ms)      — pausa sem soltar locks",
    "join()         — espera outra thread terminar",
    "yield()        — sugere ceder o processador",
    "interrupt()    — pede para parar (cooperativo)",
    "isInterrupted()— consulta a flag",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=22)

# 5 — Ciclo de vida (grafico)
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "3) Ciclo de vida da Thread")
s.shapes.add_picture(grafico_ciclo_vida(), Inches(1.5), Inches(1.4), width=Inches(10.5))

# 6 — Daemon vs User
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "3) Daemon vs User Threads")
add_bullets(s, [
    "User thread: a JVM espera ela terminar",
    "Daemon thread: morre junto com a JVM (ex.: GC)",
    "setDaemon(true) ANTES de start()",
    "Quando só restam daemons → JVM encerra",
], Inches(0.8), Inches(1.4), Inches(12), Inches(4), size=22)
add_code(s, """Thread t = new Thread(...);
t.setDaemon(true);
t.start();""", Inches(0.8), Inches(5.0), Inches(8), Inches(1.8), size=16)

# 7 — synchronized
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "4) synchronized")
add_bullets(s, [
    "Exclusão mútua: 1 thread por vez no bloco",
    "Garante visibilidade (publica escritas)",
    "Lock implícito = this (método) ou objeto (bloco)",
    "Reentrante: mesma thread pode re-adquirir",
], Inches(0.6), Inches(1.4), Inches(6.2), Inches(4))
add_code(s, """public synchronized
void incrementar() {
    valor++;
}""", Inches(7), Inches(1.6), Inches(5.5), Inches(2.4), size=16)

# 8 — volatile
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "4) volatile")
add_bullets(s, [
    "Visibilidade entre threads (não fica em cache local)",
    "NÃO garante atomicidade (i++ continua quebrado)",
    "Uso clássico: flag de parada",
    "Mais leve que synchronized — quando basta ler/escrever",
], Inches(0.6), Inches(1.4), Inches(6.2), Inches(4))
add_code(s, """private static
volatile boolean rodando = true;

while (rodando) { ... }
// outra thread: rodando = false;""", Inches(7), Inches(1.6), Inches(5.7), Inches(3), size=15)

# 9 — final / imutabilidade
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "4) final & Imutabilidade")
add_bullets(s, [
    "final = atribuído uma vez (no construtor)",
    "Objetos imutáveis são thread-safe por natureza",
    "Safe publication: após o construtor, finals visíveis",
    "Operações retornam NOVOS objetos (copy)",
], Inches(0.8), Inches(1.4), Inches(12), Inches(4), size=22)

# 10 — Monitores
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "5) Monitores (intrinsic lock)")
add_bullets(s, [
    "Todo objeto tem um monitor",
    "synchronized(obj) adquire o monitor de obj",
    "Lock REENTRANTE: mesma thread re-entra",
    "Base de wait/notify (precisa estar com o monitor)",
], Inches(0.8), Inches(1.4), Inches(12), Inches(4), size=22)

# 11 — wait / notify
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "5) wait / notify / notifyAll")
add_bullets(s, [
    "wait()      libera o monitor e bloqueia",
    "notify()    acorda UMA thread esperando",
    "notifyAll() acorda TODAS",
    "Sempre dentro de synchronized e em while (spurious wakeups)",
    "Caso clássico: produtor-consumidor",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=22)

# 12 — JMM happens-before
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "6) Java Memory Model — happens-before")
add_bullets(s, [
    "Se A happens-before B, efeitos de A são visíveis em B",
    "start() de uma thread happens-before seu início",
    "Tudo da thread happens-before o join() retornar",
    "unlock happens-before o próximo lock do mesmo monitor",
    "write em volatile happens-before read seguinte",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=20)

# 13 — Problemas (grafico)
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "7) Problemas Clássicos")
s.shapes.add_picture(grafico_problemas(), Inches(0.6), Inches(1.4), width=Inches(7))
add_bullets(s, [
    "Race: resultado depende da ordem",
    "Deadlock: locks em ordem inversa",
    "Starvation: thread nunca pega CPU/lock",
    "Livelock: reagem entre si, não progridem",
], Inches(8), Inches(1.6), Inches(5), Inches(5), size=18)

# 14 — Race condition (grafico)
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "7) Race Condition — medida")
s.shapes.add_picture(grafico_race_condition(), Inches(1.2), Inches(1.4), width=Inches(11))

# 15 — Deadlock
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "7) Deadlock")
add_bullets(s, [
    "Duas threads, dois locks, ordem invertida → trava tudo",
    "Como evitar: ordem global de aquisição",
    "Alternativa: ReentrantLock.tryLock(timeout)",
], Inches(0.6), Inches(1.4), Inches(12), Inches(2.5), size=20)
add_code(s, """// t1                       // t2
synchronized(A) {           synchronized(B) {
  synchronized(B) {           synchronized(A) {
    ...                          ...   // DEADLOCK
  }                           }
}                           }""", Inches(0.6), Inches(4.2), Inches(12), Inches(2.5), size=14)

# 16 — Starvation & Livelock
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "7) Starvation & Livelock")
add_bullets(s, [
    "Starvation: prioridades extremas ou thread gulosa monopoliza o lock",
    "Livelock: cedem educadamente pra sempre — sem deadlock, sem progresso",
    "Solução: fairness (ReentrantLock(true)), backoff aleatório",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=22)

# 17 — ThreadLocal
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "8) ThreadLocal")
add_bullets(s, [
    "Cada thread tem sua PRÓPRIA cópia",
    "Útil para: requestId, SimpleDateFormat, contexto de usuário",
    "Em thread pool: chamar remove() para evitar leak",
], Inches(0.6), Inches(1.4), Inches(6.2), Inches(4))
add_code(s, """ThreadLocal<Integer> c =
  ThreadLocal.withInitial(() -> 0);
c.set(c.get() + 1);
c.remove();""", Inches(7), Inches(1.6), Inches(5.7), Inches(3), size=15)

# 18 — InheritableThreadLocal
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "8) InheritableThreadLocal")
add_bullets(s, [
    "Igual ao ThreadLocal, mas o valor é HERDADO pela thread-filha",
    "Snapshot tirado no momento da criação da filha",
    "Útil p/ propagar contexto (requestId, locale, tracing)",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=22)

# 19 — Trade-off (grafico)
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "Trade-off entre técnicas")
s.shapes.add_picture(grafico_economia(), Inches(1.2), Inches(1.4), width=Inches(11))

# 20 — Demo
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "Demo do Código")
add_bullets(s, [
    "16 exercícios independentes na pasta src/",
    "Cada arquivo tem main + comentário-resumo no topo",
    "Compilar:  javac -d out src/*.java",
    "Executar:  java -cp out Ex11_RaceCondition",
    "Rodaremos ao vivo: Ex11 (race), Ex12 (deadlock), Ex09 (wait/notify)",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=22)

# 21 — Conclusao
s = add_slide(prs); set_background(s, FUNDO); add_title_bar(s, "Conclusão")
add_bullets(s, [
    "Concorrência ≠ paralelismo — é coordenação",
    "Prefira imutabilidade > volatile > synchronized > wait/notify",
    "Sempre adquira locks na mesma ordem",
    "ThreadLocal isola estado por thread sem locks",
    "Entender o JMM evita bugs invisíveis em produção",
], Inches(0.8), Inches(1.4), Inches(12), Inches(5), size=22)

# 22 — Obrigado
s = add_slide(prs); set_background(s, FUNDO)
add_text(s, "Obrigado!", Inches(0.6), Inches(2.5), Inches(12), Inches(1.5),
         size=80, bold=True, color=AZUL)
add_text(s, "Perguntas?", Inches(0.6), Inches(4.0), Inches(12), Inches(1),
         size=40, color=LARANJA)
add_text(s, equipe, Inches(0.6), Inches(6.2), Inches(12), Inches(0.6), size=16, color=CINZA)

out = os.path.join(SLIDES_DIR, "Seminario_Concorrencia_Java.pptx")
prs.save(out)
print("OK ->", out)
